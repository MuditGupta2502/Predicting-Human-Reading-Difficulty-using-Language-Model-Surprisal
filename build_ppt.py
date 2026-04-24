"""
Build an expanded, detailed project PPT:
Predicting Human Reading Difficulty using Language Model Surprisal
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

# ── Paths ──────────────────────────────────────────────────────────────────────
BASE   = os.path.dirname(os.path.abspath(__file__))
FIGS   = os.path.join(BASE, "results", "report_figures")
FIGS2  = os.path.join(BASE, "results")   # for preprocessing figs saved directly
OUT    = os.path.join(BASE, "results", "Project_Presentation.pptx")

# ── Colors ─────────────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # very dark navy
ACCENT_BLUE  = RGBColor(0x1E, 0x90, 0xFF)   # dodger blue
ACCENT_TEAL  = RGBColor(0x00, 0xC9, 0xA7)   # teal
ACCENT_CORAL = RGBColor(0xFF, 0x6B, 0x6B)   # coral
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xB0, 0xBE, 0xC5)
YELLOW       = RGBColor(0xFF, 0xD7, 0x00)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ══════════════════════════════════════════════════════════════════════════════
# Helper functions
# ══════════════════════════════════════════════════════════════════════════════

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    fill_bg(slide, DARK_BG)
    return slide

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width     = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    return txb

def add_multiline(slide, lines, l, t, w, h,
                  font_size=Pt(16), color=WHITE,
                  bold_first=False, line_spacing=Pt(4),
                  bullet=False):
    """lines: list of (text, bold, color_override, size_override) or just strings"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, bld, col, sz = item, False, color, font_size
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            col = item[2] if len(item) > 2 else color
            sz  = item[3] if len(item) > 3 else font_size

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.space_before = line_spacing
        run = p.add_run()
        pfx = "• " if bullet and txt.strip() else ""
        run.text = pfx + txt
        run.font.size  = sz
        run.font.bold  = bld
        run.font.color.rgb = col
    return txb

def add_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        print(f"  [WARN] image not found: {path}")
        return
    if h:
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        slide.shapes.add_picture(path, l, t, w)

def accent_bar(slide, color=ACCENT_BLUE, height=Inches(0.06)):
    """Thin colored bar across top of slide."""
    add_rect(slide, 0, 0, SLIDE_W, height, fill_color=color)

def section_label(slide, text, color=ACCENT_TEAL):
    add_text(slide, text,
             Inches(0.4), Inches(0.08), Inches(6), Inches(0.35),
             font_size=Pt(11), bold=True, color=color, align=PP_ALIGN.LEFT)

def slide_title(slide, text, sub=None, title_color=WHITE):
    accent_bar(slide, ACCENT_BLUE)
    add_text(slide, text,
             Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.7),
             font_size=Pt(32), bold=True, color=title_color, align=PP_ALIGN.LEFT)
    if sub:
        add_text(slide, sub,
                 Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
                 font_size=Pt(17), bold=False, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

def divider(slide, y, color=ACCENT_BLUE, alpha=0.4):
    add_rect(slide, Inches(0.5), y, Inches(12.33), Inches(0.025),
             fill_color=color)

def card(slide, l, t, w, h, fill=RGBColor(0x15, 0x2A, 0x40), border=ACCENT_BLUE):
    add_rect(slide, l, t, w, h, fill_color=fill, line_color=border, line_width=Pt(1.2))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
# gradient-like overlay
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill_color=DARK_BG)
add_rect(s, 0, 0, Inches(4.5), SLIDE_H, fill_color=RGBColor(0x05, 0x10, 0x20))
# accent strip
add_rect(s, Inches(4.4), 0, Inches(0.08), SLIDE_H, fill_color=ACCENT_BLUE)

add_text(s, "Predicting Human Reading Difficulty",
         Inches(4.7), Inches(1.6), Inches(8.2), Inches(1.0),
         font_size=Pt(34), bold=True, color=WHITE)
add_text(s, "using Language Model Surprisal",
         Inches(4.7), Inches(2.55), Inches(8.2), Inches(0.7),
         font_size=Pt(28), bold=False, color=ACCENT_TEAL)
add_text(s, "Computational Psycholinguistics — Course Project",
         Inches(4.7), Inches(3.4), Inches(8.2), Inches(0.4),
         font_size=Pt(14), color=LIGHT_GREY)
add_text(s, "Mudit Gupta  |  Roll: 2024201058",
         Inches(4.7), Inches(3.85), Inches(8.2), Inches(0.4),
         font_size=Pt(13), color=LIGHT_GREY)

add_text(s, "AGENDA",
         Inches(0.3), Inches(1.3), Inches(3.8), Inches(0.4),
         font_size=Pt(11), bold=True, color=ACCENT_BLUE)
toc = [
    "1. Motivation & Surprisal Theory",
    "2. The Cognitive Phenomena",
    "3. Language Models as Proxies",
    "4. The Subword Alignment Problem",
    "5. Measuring Reading Difficulty",
    "6. Formulation & Hypotheses",
    "7. Datasets & Preprocessing (Nb 1)",
    "8. Classical N-gram Surprisal (Nb 2)",
    "9. Transformer Context Limits (Nb 3)",
    "10. Main Findings & Visuals",
    "11. Why did BERT fail?",
    "12. Future Directions",
]
add_multiline(s, toc, Inches(0.3), Inches(1.75), Inches(3.9), Inches(5.0),
              font_size=Pt(12), color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Background & Motivation
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_TEAL)
section_label(s, "BACKGROUND", ACCENT_TEAL)
slide_title(s, "Background & Motivation", "Cognitive load is not distributed evenly across a sentence")

add_text(s, "Why do some words take longer to read than others?",
         Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         font_size=Pt(20), bold=True, color=ACCENT_TEAL)

add_text(s, "Core concept: Surprisal Theory   (Hale, 2001; Levy, 2008)",
         Inches(0.5), Inches(1.95), Inches(9), Inches(0.4),
         font_size=Pt(15), bold=True, color=YELLOW)

add_text(s, "S(wᵢ) = −log₂ P(wᵢ | context)",
         Inches(0.7), Inches(2.4), Inches(8), Inches(0.5),
         font_size=Pt(19), bold=True, color=ACCENT_BLUE)

bullets = [
    ("Information Theory meets Cognition:", True, WHITE, Pt(15)),
    ("A word that is highly predictable carries little information, thus requiring little processing effort.", False, LIGHT_GREY, Pt(14)),
    ("Conversely, an unexpected word (high surprisal) forces the brain to update its syntactic/semantic trace rapidly, causing a delay.", False, LIGHT_GREY, Pt(14)),
    ("", False, LIGHT_GREY, Pt(14)),
    ("The Question:", True, WHITE, Pt(15)),
    ("Can artificial neural networks predict human cognitive load?", False, ACCENT_TEAL, Pt(14)),
    ("Do large language models (which implicitly learn syntax and semantics) align better with human reading behavior than classical statistical approaches?", False, LIGHT_GREY, Pt(14)),
]
add_multiline(s, bullets, Inches(0.7), Inches(3.0), Inches(8), Inches(4.0),
              bullet=True)

# Right panel — formula box
card(s, Inches(9.3), Inches(1.3), Inches(3.6), Inches(5.0))
add_text(s, "Key Mechanisms",
         Inches(9.4), Inches(1.4), Inches(3.4), Inches(0.35),
         font_size=Pt(12), bold=True, color=ACCENT_BLUE)
right_items = [
    "Predictive Processing",
    "The brain constantly guesses the next word.",
    "",
    "Cost of Update",
    "When predictions fail, reading time spikes.",
    "",
    "Logarithmic Scaling",
    "Processing difficulty scales linearly with the negative log probability (bits).",
    "",
    "RT = α * S(w) + β"
]
add_multiline(s, right_items, Inches(9.4), Inches(1.8), Inches(3.3), Inches(4.0),
              font_size=Pt(13), color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Cognitive Phenomena (Garden Path)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_TEAL)
section_label(s, "COGNITIVE LINGUISTICS", ACCENT_TEAL)
slide_title(s, "The Target Phenomena: Syntactic Complexity", "What specifically causes surprisal spikes?")

card(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(2.5))
add_text(s, "1. Garden-Path Sentences", Inches(0.6), Inches(1.5), Inches(5.0), Inches(0.4), font_size=Pt(16), bold=True, color=ACCENT_CORAL)
gp_text = [
    ("The complex houses married and single soldiers.", True, WHITE, Pt(16)),
    ("At first glance, 'complex' is read as an adjective and 'houses' as a noun.", False, LIGHT_GREY, Pt(14)),
    ("By 'married', the parse fails. Regression happens. 'complex' is actually a noun, 'houses' is a verb.", False, LIGHT_GREY, Pt(14)),
    ("→ Massive reading time spike at 'married'. Models should predict high surprisal here.", False, ACCENT_TEAL, Pt(14)),
]
add_multiline(s, gp_text, Inches(0.6), Inches(1.9), Inches(12.0), Inches(1.8), bullet=True)


card(s, Inches(0.4), Inches(4.1), Inches(12.5), Inches(2.6))
add_text(s, "2. Center-Embedding", Inches(0.6), Inches(4.2), Inches(5.0), Inches(0.4), font_size=Pt(16), bold=True, color=ACCENT_CORAL)
ce_text = [
    ("The rat the cat the dog chased killed ate the malt.", True, WHITE, Pt(16)),
    ("Working memory is overwhelmed maintaining multiple open syntactic dependencies.", False, LIGHT_GREY, Pt(14)),
    ("Reading times increase linearly with embedding depth.", False, LIGHT_GREY, Pt(14)),
    ("→ Does a Transformer's self-attention mechanism replicate this memory strain via surprisal?", False, ACCENT_TEAL, Pt(14)),
]
add_multiline(s, ce_text, Inches(0.6), Inches(4.6), Inches(12.0), Inches(1.8), bullet=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Prerequisites (LMs)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CORAL)
section_label(s, "PREREQUISITES", ACCENT_CORAL)
slide_title(s, "Prerequisites — Language Models", "Three distinct architectures for estimating P(w|context)")

# N-gram card
card(s, Inches(0.4), Inches(1.4), Inches(4.0), Inches(5.5))
add_text(s, "📊  N-gram (KN)",
         Inches(0.55), Inches(1.5), Inches(3.6), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT_CORAL)
ngram_pts = [
    ("Predicts next word using ONLY the last N−1 words.", False, LIGHT_GREY, Pt(13)),
    ("P(wᵢ | wᵢ₋₁, wᵢ₋₂)", False, ACCENT_BLUE, Pt(13)),
    ("", False, WHITE, Pt(8)),
    ("Kneser-Ney Smoothing:", True, YELLOW, Pt(13)),
    ("Handles unseen n-grams gracefully by backing off to lower orders based on continuation probabilities.", False, LIGHT_GREY, Pt(12)),
    ("", False, WHITE, Pt(8)),
    ("Strengths:", True, WHITE, Pt(13)),
    ("Extremely robust locally.", False, LIGHT_GREY, Pt(12)),
    ("Weakness:", True, WHITE, Pt(13)),
    ("Has no concept of long-range syntax (subject-verb agreement over distance).", False, LIGHT_GREY, Pt(12)),
]
add_multiline(s, ngram_pts, Inches(0.55), Inches(1.95), Inches(3.6), Inches(4.5))

# GPT-2 card
card(s, Inches(4.6), Inches(1.4), Inches(4.0), Inches(5.5))
add_text(s, "⚡ GPT-2 (Autoregressive)",
         Inches(4.75), Inches(1.5), Inches(3.8), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT_TEAL)
transf_pts = [
    ("Causal Transformer.", True, YELLOW, Pt(13)),
    ("Predicts each token using dense representations of ALL prior tokens via self-attention.", False, LIGHT_GREY, Pt(12)),
    ("P(wᵢ | w₁...wᵢ₋₁)", False, ACCENT_BLUE, Pt(12)),
    ("", False, WHITE, Pt(6)),
    ("Strengths:", True, WHITE, Pt(13)),
    ("Context window of 1024 tokens. Can track characters and complex grammar.", False, LIGHT_GREY, Pt(12)),
    ("Close analog to human sequential reading.", False, LIGHT_GREY, Pt(12)),
    ("Weakness:", True, WHITE, Pt(13)),
    ("Only single-direction context.", False, LIGHT_GREY, Pt(12)),
]
add_multiline(s, transf_pts, Inches(4.75), Inches(1.95), Inches(3.7), Inches(4.5))

# BERT card
card(s, Inches(8.8), Inches(1.4), Inches(3.9), Inches(5.5))
add_text(s, "🎭 BERT (Masked LM)",
         Inches(8.95), Inches(1.5), Inches(3.7), Inches(0.4),
         font_size=Pt(15), bold=True, color=YELLOW)
bert_pts = [
    ("Bidirectional Transformer.", True, YELLOW, Pt(13)),
    ("Not designed for left-to-right prediction. Uses Pseudo-Log-Likelihood (PLL).", False, LIGHT_GREY, Pt(12)),
    ("Mask target word, then read log P from [MASK] logit.", False, ACCENT_BLUE, Pt(12)),
    ("", False, WHITE, Pt(6)),
    ("Strengths:", True, WHITE, Pt(13)),
    ("Sees future context. Superior for semantic embedding.", False, LIGHT_GREY, Pt(12)),
    ("Weakness:", True, WHITE, Pt(13)),
    ("Not physiologically plausible (humans don't see the end of the sentence first).", False, LIGHT_GREY, Pt(12)),
    ("Computationally heavy (1 pass per word).", False, LIGHT_GREY, Pt(12)),
]
add_multiline(s, bert_pts, Inches(8.95), Inches(1.95), Inches(3.7), Inches(4.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Alignment Problem
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CORAL)
section_label(s, "METHODOLOGY CHALLENGES", ACCENT_CORAL)
slide_title(s, "The Subword Alignment Problem", "Reconciling model tokenizers with human reading data")

card(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.5))

align_text = [
    ("The Mismatch:", True, YELLOW, Pt(17)),
    ("Psycholinguistic corpora (like Natural Stories) record reading times per WORD.", False, LIGHT_GREY, Pt(15)),
    ("Modern transformers (GPT-2, BERT) use Byte Pair Encoding (BPE) or WordPiece.", False, LIGHT_GREY, Pt(15)),
    ("A single word is often split into a sequence of subword tokens.", False, LIGHT_GREY, Pt(15)),
    ("", False, WHITE, Pt(10)),
    ("Example:", True, ACCENT_TEAL, Pt(16)),
    ("Corpus word:  'unbelievable'   →   RT: 450 ms", False, WHITE, Pt(15)),
    ("GPT-2 output: ['un', 'bel', 'iev', 'able']   →   We have 4 separate log probabilities.", False, WHITE, Pt(15)),
    ("", False, WHITE, Pt(10)),
    ("Mathematical Resolution (Joint Probability):", True, YELLOW, Pt(17)),
    ("The probability of the full word is the joint probability of its subwords.", False, LIGHT_GREY, Pt(15)),
    ("Because probabilities multiply, log-probabilities ADD.", False, LIGHT_GREY, Pt(15)),
    ("S(word) = −( log₂ P('un' | ctx) + log₂ P('bel' | ctx + 'un') + ... )", True, ACCENT_BLUE, Pt(16)),
    ("", False, WHITE, Pt(10)),
    ("Implementation detail:", True, ACCENT_CORAL, Pt(16)),
    ("We must carefully track token indices back to the original text sequence to ensure every subword's log-prob is rolled up accurately into the final CSVs.", False, LIGHT_GREY, Pt(14)),
]
add_multiline(s, align_text, Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Measuring Reading Difficulty
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CORAL)
section_label(s, "MEASUREMENT", ACCENT_CORAL)
slide_title(s, "Measuring Human Reading Difficulty", "Self-paced reading vs Eye-tracking")

# Left: SPR
card(s, Inches(0.4), Inches(1.4), Inches(6.1), Inches(5.5))
add_text(s, "📖  Self-Paced Reading (SPR)",
         Inches(0.55), Inches(1.5), Inches(5.7), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT_TEAL)
rt_pts = [
    ("Mechanism:", True, YELLOW, Pt(13)),
    ("Participant sits at a computer. Text is hidden by dashes ('--- ---').", False, LIGHT_GREY, Pt(13)),
    ("Pressing spacebar reveals next word and hides previous word.", False, LIGHT_GREY, Pt(13)),
    ("RT = time between spacebar presses (ms).", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(8)),
    ("Pros:", True, WHITE, Pt(13)),
    ("Forces strict serial processing. Cannot look back.", False, LIGHT_GREY, Pt(13)),
    ("Very clear measure of localized difficulty.", False, LIGHT_GREY, Pt(13)),
    ("Cons:", True, WHITE, Pt(13)),
    ("Unnatural reading format. Slower than normal reading.", False, LIGHT_GREY, Pt(13)),
    ("Spillover effects common (delayed reaction on word i+1).", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(8)),
    ("Used in: Natural Stories Corpus", True, ACCENT_BLUE, Pt(13)),
]
add_multiline(s, rt_pts, Inches(0.55), Inches(1.95), Inches(5.5), Inches(4.5))

# Right: Eye-tracking
card(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.5))
add_text(s, "👁️  Eye-Tracking",
         Inches(7.05), Inches(1.5), Inches(5.6), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT_CORAL)
et_pts = [
    ("Mechanism:", True, YELLOW, Pt(13)),
    ("Infrared camera tracks pupil/corneal reflection at 1000 Hz.", False, LIGHT_GREY, Pt(13)),
    ("Records 'saccades' (jumps) and 'fixations' (resting).", False, LIGHT_GREY, Pt(13)),
    ("We use First-Pass Fixation Duration (FDUR).", False, ACCENT_BLUE, Pt(13)),
    ("", False, WHITE, Pt(8)),
    ("Pros:", True, WHITE, Pt(13)),
    ("Highly naturalistic. Participants read paragraphs normally.", False, LIGHT_GREY, Pt(13)),
    ("Allows studying regressions (looking back).", False, LIGHT_GREY, Pt(13)),
    ("Cons:", True, WHITE, Pt(13)),
    ("Word skipping (short words often have 0 ms fixation).", False, LIGHT_GREY, Pt(13)),
    ("Highly complex data requiring parsing.", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(8)),
    ("Used in: Dundee Corpus", True, ACCENT_BLUE, Pt(13)),
]
add_multiline(s, et_pts, Inches(7.05), Inches(1.95), Inches(5.6), Inches(4.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Research Questions
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_BLUE)
section_label(s, "RESEARCH QUESTIONS", ACCENT_BLUE)
slide_title(s, "Research Questions & Hypotheses")

rqs = [
    ("RQ1", "Does language-model surprisal correlate with human reading time?",
     "Higher surprisal → longer reading time (RT ∝ Surprisal). We expect positive r.", ACCENT_BLUE),
    ("RQ2", "Do transformer LMs (GPT-2/BERT) predict reading difficulty better than n-gram LMs?",
     "Transformers capture long-range context → higher correlation with RT than Bigrams.", ACCENT_TEAL),
    ("RQ3", "Do syntactic constructions produce surprisal spikes matching RT increases?",
     "Garden-path & center-embedded sentences → local surprisal peaks at disambiguation points.", ACCENT_CORAL),
]

y = Inches(1.4)
for rq_id, question, hyp, col in rqs:
    card(s, Inches(0.4), y, Inches(12.5), Inches(1.6), border=col)
    add_text(s, rq_id, Inches(0.5), y + Inches(0.12), Inches(0.85), Inches(0.45),
             font_size=Pt(20), bold=True, color=col)
    add_text(s, question, Inches(1.4), y + Inches(0.1), Inches(11.1), Inches(0.55),
             font_size=Pt(15), bold=True, color=WHITE)
    add_text(s, "Hypothesis: " + hyp, Inches(1.4), y + Inches(0.72), Inches(11.1), Inches(0.6),
             font_size=Pt(13), bold=False, color=LIGHT_GREY, italic=True)
    y += Inches(1.75)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Datasets
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_TEAL)
section_label(s, "DATASETS", ACCENT_TEAL)
slide_title(s, "Datasets Used in This Project")

# NS card
card(s, Inches(0.4), Inches(1.4), Inches(6.0), Inches(5.6))
add_text(s, "📗  Natural Stories Corpus",
         Inches(0.55), Inches(1.5), Inches(5.6), Inches(0.4),
         font_size=Pt(16), bold=True, color=ACCENT_TEAL)

ns_pts = [
    ("Futrell et al. (2018)", True, YELLOW, Pt(13)),
    ("", False, WHITE, Pt(5)),
    ("• 10 natural language stories", False, LIGHT_GREY, Pt(13)),
    ("• 180 participants", False, LIGHT_GREY, Pt(13)),
    ("• 10,256 unique word tokens", False, LIGHT_GREY, Pt(13)),
    ("• Designed to contain complex syntactic", False, LIGHT_GREY, Pt(13)),
    ("  structures (garden-path, embeddings)", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(5)),
    ("Cleaning steps:", True, YELLOW, Pt(13)),
    ("• Limit RT to 100–3000 ms to filter out slips", False, LIGHT_GREY, Pt(12)),
    ("• The Z-Score Imperative:", True, ACCENT_CORAL, Pt(12)),
    ("   Some people read fast, some slow.", False, LIGHT_GREY, Pt(12)),
    ("   We z-score log(RT) per participant to isolate", False, LIGHT_GREY, Pt(12)),
    ("   relative difficulty of specific words.", False, LIGHT_GREY, Pt(12)),
    ("", False, WHITE, Pt(5)),
    ("DV: mean log(RT) per word", True, ACCENT_BLUE, Pt(13)),
]
add_multiline(s, ns_pts, Inches(0.55), Inches(1.95), Inches(5.6), Inches(4.8))

# Dundee card
card(s, Inches(6.8), Inches(1.4), Inches(6.1), Inches(5.6))
add_text(s, "📘  Dundee Corpus",
         Inches(6.95), Inches(1.5), Inches(5.6), Inches(0.4),
         font_size=Pt(16), bold=True, color=ACCENT_BLUE)

du_pts = [
    ("Kennedy et al. (2003)", True, YELLOW, Pt(13)),
    ("", False, WHITE, Pt(5)),
    ("• 20 UK newspaper article texts", False, LIGHT_GREY, Pt(13)),
    ("• 10 participants (but dense data)", False, LIGHT_GREY, Pt(13)),
    ("• 50,650 word-level observations", False, LIGHT_GREY, Pt(13)),
    ("• General domain text (not crafted)", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(5)),
    ("Cleaning steps:", True, YELLOW, Pt(13)),
    ("• Keep Pass-1 fixations only. Why?", True, ACCENT_CORAL, Pt(12)),
    ("   Re-reading (regression) signifies failure of", False, LIGHT_GREY, Pt(12)),
    ("   surprisal predictions. First-pass isolates", False, LIGHT_GREY, Pt(12)),
    ("   local predictability.", False, LIGHT_GREY, Pt(12)),
    ("• FDUR bounded 50–1200 ms", False, LIGHT_GREY, Pt(12)),
    ("• Z-score log(FDUR) per subject", False, LIGHT_GREY, Pt(12)),
    ("", False, WHITE, Pt(5)),
    ("DV: mean log(FDUR) per word", True, ACCENT_BLUE, Pt(13)),
]
add_multiline(s, du_pts, Inches(6.95), Inches(1.95), Inches(5.6), Inches(4.8))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Notebook 1 Figures
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_TEAL)
section_label(s, "NOTEBOOK 1 — PREPROCESSING", ACCENT_TEAL)
slide_title(s, "Notebook 1: Distributions & Skew Rectification")

add_image(s, os.path.join(FIGS2, "report_figures", "fig_ns_rt.png"),
          Inches(0.3), Inches(1.3), Inches(4.2))
add_text(s, "Natural Stories RT", Inches(0.3), Inches(4.4),
         Inches(4.2), Inches(0.35), font_size=Pt(12), color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

add_image(s, os.path.join(FIGS2, "report_figures", "fig_du_fdur.png"),
          Inches(4.6), Inches(1.3), Inches(4.2))
add_text(s, "Dundee Fixation Duration", Inches(4.6), Inches(4.4),
         Inches(4.2), Inches(0.35), font_size=Pt(12), color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

add_image(s, os.path.join(FIGS2, "report_figures", "fig_overlay.png"),
          Inches(8.9), Inches(1.3), Inches(4.0))
add_text(s, "Log-DV Overlay", Inches(8.9), Inches(4.4),
         Inches(4.0), Inches(0.35), font_size=Pt(12), color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_multiline(s, [
    ("Right-Skew Issue:", True, YELLOW, Pt(14)),
    ("Reaction times naturally have a strict lower bound (0 ms) but an infinite upper tail.", False, LIGHT_GREY, Pt(13)),
    ("OLS Regression assumes normality. Without a logarithmic transform, outliers pull the trend line.", False, LIGHT_GREY, Pt(13)),
    ("Log-transformation normalizes the data as shown in the orange distributions.", False, LIGHT_GREY, Pt(13)),
], Inches(0.4), Inches(5.0), Inches(12.8), Inches(2.1))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Notebook 2: N-gram Surprisal
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_BLUE)
section_label(s, "NOTEBOOK 2", ACCENT_BLUE)
slide_title(s, "Notebook 2: Classical N-gram Surprisal", "Establishing the baseline with Kneser-Ney smoothing")

# Left: Method
card(s, Inches(0.4), Inches(1.4), Inches(6.2), Inches(5.5))
add_text(s, "Method", Inches(0.55), Inches(1.5), Inches(5.8), Inches(0.4),
         font_size=Pt(16), bold=True, color=ACCENT_BLUE)
method_pts = [
    ("Training corpus:", True, YELLOW, Pt(13)),
    ("NLTK Brown corpus — 57,340 sentences. Generic domain.", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(5)),
    ("Models trained:", True, YELLOW, Pt(13)),
    ("• Bigram  (n=2)", False, ACCENT_BLUE, Pt(13)),
    ("• Trigram (n=3)", False, ACCENT_BLUE, Pt(13)),
    ("", False, WHITE, Pt(5)),
    ("Implementation:", True, YELLOW, Pt(13)),
    ("NLTK's `KneserNeyInterpolated` pipeline.", False, LIGHT_GREY, Pt(13)),
    ("Computed sequentially across texts. OOV tokens smoothed.", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(5)),
    ("Output generated:", True, YELLOW, Pt(13)),
    ("ns_ngram_surprisal.csv (10k rows)", False, LIGHT_GREY, Pt(12)),
    ("dundee_ngram_surprisal.csv (50k rows)", False, LIGHT_GREY, Pt(12)),
]
add_multiline(s, method_pts, Inches(0.55), Inches(1.95), Inches(5.8), Inches(4.6))

# Right: Results
card(s, Inches(6.85), Inches(1.4), Inches(6.1), Inches(5.5))
add_text(s, "Results", Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.4),
         font_size=Pt(16), bold=True, color=ACCENT_CORAL)

add_text(s, "Baseline Correlation with log-DV:",
         Inches(7.0), Inches(2.2), Inches(5.7), Inches(0.3),
         font_size=Pt(13), bold=True, color=YELLOW)

# Correlation table
table_data = [
    ["Model",        "NS r(log RT)", "Dundee r(logFDUR)"],
    ["Bigram KN",  "r = 0.2149",   "r = 0.1732"],
    ["Trigram KN", "r = 0.1928",   "r = 0.1715"],
]
row_colors  = [ACCENT_BLUE, RGBColor(0x15, 0x2A, 0x40), RGBColor(0x0F, 0x20, 0x30)]
col_widths  = [Inches(2.0), Inches(1.9), Inches(2.1)]
ty = Inches(2.7)
for ri, row in enumerate(table_data):
    tx = Inches(7.0)
    for ci, cell_txt in enumerate(row):
        bg = ACCENT_BLUE if ri == 0 else (RGBColor(0x15, 0x2A, 0x40) if ri % 2 == 1 else RGBColor(0x0F, 0x20, 0x30))
        add_rect(s, tx, ty, col_widths[ci], Inches(0.4), fill_color=bg)
        add_text(s, cell_txt, tx + Inches(0.05), ty + Inches(0.05),
                 col_widths[ci] - Inches(0.1), Inches(0.32),
                 font_size=Pt(12 if ri > 0 else 11),
                 bold=(ri == 0), color=WHITE if ri == 0 else LIGHT_GREY,
                 align=PP_ALIGN.CENTER)
        tx += col_widths[ci]
    ty += Inches(0.42)

add_text(s, "Why did Bigram beat Trigram?",
         Inches(7.0), Inches(4.5), Inches(5.7), Inches(0.35),
         font_size=Pt(12), bold=True, color=ACCENT_CORAL)
add_multiline(s, ["Sparsity. Trigrams suffer more from data sparsity on a small training corpus (Brown), introducing noise."], 
              Inches(7.0), Inches(4.9), Inches(5.7), Inches(1.0), font_size=Pt(12), color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Notebook 3: Transformer Engineering Hurdles
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CORAL)
section_label(s, "NOTEBOOK 3 — FIXES", ACCENT_CORAL)
slide_title(s, "Notebook 3: Transformer Context Limits", "Engineering around CUDA / Memory crashes")

card(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.5))

fix_text = [
    ("The Problem:", True, ACCENT_CORAL, Pt(17)),
    ("BERT max sequence length is 512 subwords.", False, LIGHT_GREY, Pt(15)),
    ("GPT-2 max sequence length is 1024 subwords.", False, LIGHT_GREY, Pt(15)),
    ("Dundee articles are often 2000+ words. Passing the whole text yields a CUDA Device Assert on indexing.", False, LIGHT_GREY, Pt(15)),
    ("", False, WHITE, Pt(10)),
    ("GPT-2 Solution (Autoregressive Context Tracking):", True, ACCENT_BLUE, Pt(16)),
    ("We built a chunker that grabs words greedily until subword tokens hit 512.", False, WHITE, Pt(15)),
    ("Only the first word of the chunk loses left context. This minimizes data loss across the vast sequence.", False, WHITE, Pt(15)),
    ("", False, WHITE, Pt(10)),
    ("BERT Solution (PLL Tracking):", True, YELLOW, Pt(16)),
    ("Created 80-word windows. In each window, loop through every word, MASK it, do a forward pass, and sum", False, LIGHT_GREY, Pt(15)),
    ("local subword logits. Extends runtime significantly but avoids out-of-bounds truncation.", False, LIGHT_GREY, Pt(15)),
]
add_multiline(s, fix_text, Inches(0.6), Inches(1.6), Inches(12.0), Inches(5.0))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Main Results: Correlation Table
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_BLUE)
section_label(s, "RESULTS", ACCENT_BLUE)
slide_title(s, "Main Results: 4-Model Correlation Table")

add_text(s, "Pearson r between model surprisal and human log reading time — all p < 0.001",
         Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
         font_size=Pt(14), color=LIGHT_GREY, italic=True)

# Big table
headers = ["Model", "Type", "NS r(log RT)", "Dundee r(logFDUR)", "Verdict"]
rows_data = [
    ["Bigram KN",  "Classical N-gram", "0.2149 ⭐", "0.1732", "Best on NS  (1st)"],
    ["Trigram KN", "Classical N-gram", "0.1928",   "0.1715", ""],
    ["GPT-2",      "Transformer (AR)", "0.2094",   "0.1929 ⭐", "Best on DU (1st)"],
    ["BERT (PLL)", "Transformer (MLM)","0.1572",   "0.1554", "Weakest (4th)"],
]
col_ws = [Inches(2.1), Inches(2.3), Inches(2.5), Inches(2.8), Inches(2.8)]
tx_start = Inches(0.4)
ty = Inches(1.85)

for ri, row in enumerate([headers] + rows_data):
    tx = tx_start
    header_row = ri == 0
    for ci, cell_txt in enumerate(row):
        if header_row:
            bg = ACCENT_BLUE
        elif ri % 2 == 1:
            bg = RGBColor(0x12, 0x24, 0x38)
        else:
            bg = RGBColor(0x0D, 0x1B, 0x2A)

        # Highlight best values
        txt_color = WHITE if header_row else LIGHT_GREY
        if "⭐" in cell_txt: txt_color = YELLOW
        if "Weakest" in cell_txt: txt_color = ACCENT_CORAL
        if "Best" in cell_txt: txt_color = ACCENT_TEAL

        add_rect(s, tx, ty, col_ws[ci], Inches(0.52), fill_color=bg)
        add_text(s, cell_txt, tx + Inches(0.07), ty + Inches(0.1),
                 col_ws[ci] - Inches(0.14), Inches(0.35),
                 font_size=Pt(13 if not header_row else 12),
                 bold=header_row, color=txt_color, align=PP_ALIGN.CENTER)
        tx += col_ws[ci]
    ty += Inches(0.54)

add_text(s, "⭐ = best in column    AR = autoregressive    MLM = masked language model",
         Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.35),
         font_size=Pt(11), color=LIGHT_GREY, italic=True)

# Interpretation boxes
interp = [
    ("RQ1 ✓", "All 4 models show significant positive correlation with reading time.\nSurprisal IS a reliable predictor of reading difficulty.", ACCENT_TEAL),
    ("RQ2 ~", "MIXED result: GPT-2 outperforms n-grams on Dundee but Bigram KN marginally outperforms GPT-2 on Natural Stories (0.215 vs 0.209).", YELLOW),
]

y = Inches(5.0)
for title, text, col in interp:
    add_rect(s, Inches(0.4), y, Inches(0.08), Inches(0.9), fill_color=col)
    add_text(s, title, Inches(0.6), y, Inches(1.6), Inches(0.35),
             font_size=Pt(14), bold=True, color=col)
    add_text(s, text, Inches(0.6), y + Inches(0.32), Inches(12.2), Inches(0.55),
             font_size=Pt(12), color=LIGHT_GREY)
    y += Inches(1.1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Why did BERT fail / Why did Bigram win NS?
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CORAL)
section_label(s, "ANALYSIS: THE ANOMALIES", ACCENT_CORAL)
slide_title(s, "Why did BERT fail? Why did Bigram win NS?")

card(s, Inches(0.4), Inches(1.4), Inches(6.1), Inches(5.5))
add_text(s, "Why did BERT score the worst?", Inches(0.55), Inches(1.5), Inches(5.7), Inches(0.4), font_size=Pt(16), bold=True, color=ACCENT_CORAL)
b_text = [
    ("Bidirectional cheating:", True, WHITE, Pt(14)),
    ("BERT PLL calculates P(w | Left Context + Right Context).", False, LIGHT_GREY, Pt(13)),
    ("Humans read strictly Left-to-Right. The right context is not available during an actual surprisal event in the brain.", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(8)),
    ("Noise introduced by Masking:", True, WHITE, Pt(14)),
    ("Using the [MASK] token alters the sentence dynamics compared to GPT-2's smooth autoregression.", False, LIGHT_GREY, Pt(13)),
]
add_multiline(s, b_text, Inches(0.55), Inches(2.0), Inches(5.7), Inches(4.5))

card(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(5.5))
add_text(s, "Why did Bigram win on Natural Stories?", Inches(7.05), Inches(1.5), Inches(5.6), Inches(0.4), font_size=Pt(16), bold=True, color=ACCENT_TEAL)
bg_text = [
    ("Local Syntactic Traps:", True, WHITE, Pt(14)),
    ("Natural Stories is filled with deliberate, tricky transitions (e.g. Noun directly adjacent to Verb where an Adjective was expected).", False, LIGHT_GREY, Pt(13)),
    ("", False, WHITE, Pt(8)),
    ("Bigrams are hypersensitive locally:", True, WHITE, Pt(14)),
    ("A Bigram model reacts disastrously to adjacent weird pairs. This harsh local reaction correlates perfectly with the human brain's sudden confusion at that exact boundary.", False, LIGHT_GREY, Pt(13)),
    ("GPT-2 buffers this by using massive context width to 'smooth out' expectations.", False, LIGHT_GREY, Pt(13)),
]
add_multiline(s, bg_text, Inches(7.05), Inches(2.0), Inches(5.6), Inches(4.5))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Results Figures (Bar)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_BLUE)
section_label(s, "RESULTS — VISUAL", ACCENT_BLUE)
slide_title(s, "Results: Correlation Bar Chart & Distributions")

add_image(s, os.path.join(FIGS, "fig_model_comparison.png"),
          Inches(0.3), Inches(1.3), Inches(6.5))
add_text(s, "Pearson r by model — both corpora",
         Inches(0.3), Inches(5.4), Inches(6.5), Inches(0.35),
         font_size=Pt(12), color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

add_image(s, os.path.join(FIGS, "fig_surp_distributions_all.png"),
          Inches(7.1), Inches(1.3), Inches(5.9))
add_text(s, "Surprisal distributions — all 4 models",
         Inches(7.1), Inches(5.4), Inches(5.9), Inches(0.35),
         font_size=Pt(12), color=ACCENT_CORAL, align=PP_ALIGN.CENTER)

add_multiline(s, [
    ("Key visual findings:", True, YELLOW, Pt(13)),
    ("• Bigram KN produces the highest NS correlation bar (steelblue, left-most pair)",  False, LIGHT_GREY, Pt(12)),
    ("• GPT-2 wins on Dundee (coral bar, 3rd pair) — extends bigram on broader domain text", False, LIGHT_GREY, Pt(12)),
    ("• BERT (PLL) scores heavily shifted left toward 0 bits — predicting masked words given surrounding context is 'too easy'.", False, LIGHT_GREY, Pt(12)),
], Inches(0.4), Inches(6.0), Inches(12.8), Inches(1.2))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Scatter Plot Grid
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_BLUE)
section_label(s, "RESULTS — SCATTER PLOTS", ACCENT_BLUE)
slide_title(s, "Surprisal vs. Log Reading Time — All Models × Corpora")

add_image(s, os.path.join(FIGS, "fig_scatter_all_models.png"),
          Inches(0.3), Inches(1.3), Inches(12.7))

add_multiline(s, [
    ("Each panel: one model × one corpus. Black line = OLS fit.  r value shown in title.", False, LIGHT_GREY, Pt(12)),
    ("Positive slopes in all 8 panels visually confirm: higher surprisal → longer reading time (RQ1 ✓)", True, ACCENT_TEAL, Pt(12)),
], Inches(0.4), Inches(6.7), Inches(12.8), Inches(0.7))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Discussion & Conclusions
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_TEAL)
section_label(s, "CONCLUSIONS", ACCENT_TEAL)
slide_title(s, "Discussion & Conclusions")

findings = [
    ("✅ RQ1 — Confirmed", "All LM surprisal measures positively correlate with human reading time (r=0.15–0.21). Higher surprisal reliably predicts longer reading times / fixation durations.", ACCENT_TEAL),
    ("⚠️ RQ2 — Mixed", "GPT-2 outperforms n-grams on Dundee but not on Natural Stories. Classical Bigram models are surprisingly competitive at detecting local structural ambiguity.", YELLOW),
    ("🔍 BERT Failure", "Bidirectional context is not physiologically valid for sequential reading measurement. Masked LM produces flawed surprisal metrics for psycholinguistics.", ACCENT_CORAL),
    ("📌 Replication", "Results replicate across both corpora (Self-Paced Reading and Eye-Tracking), proving standard linear association holds independent of tracking mechanism.", ACCENT_BLUE),
]

y = Inches(1.4)
for title, body, col in findings:
    card(s, Inches(0.4), y, Inches(12.5), Inches(1.35), border=col)
    add_rect(s, Inches(0.4), y, Inches(0.08), Inches(1.35), fill_color=col)
    add_text(s, title, Inches(0.6), y + Inches(0.06), Inches(4.0), Inches(0.38),
             font_size=Pt(14), bold=True, color=col)
    add_text(s, body, Inches(0.6), y + Inches(0.5), Inches(12.1), Inches(0.75),
             font_size=Pt(12.5), color=LIGHT_GREY)
    y += Inches(1.48)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — What's Next + References
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
accent_bar(s, ACCENT_CORAL)
section_label(s, "FUTURE WORK & REFERENCES", ACCENT_CORAL)
slide_title(s, "What's Next & References")

card(s, Inches(0.4), Inches(1.4), Inches(7.5), Inches(5.5))
add_text(s, "🔮  Remaining Work (50% → 100%)",
         Inches(0.55), Inches(1.5), Inches(7.1), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT_CORAL)
next_pts = [
    ("Notebook 4 — Mixed-Effects Regression:", True, YELLOW, Pt(13)),
    ("  Fit lmer-style models: logRT ~ surprisal + word_len +", False, LIGHT_GREY, Pt(12)),
    ("  word_freq + position + (1|subject) + (1|item)", False, LIGHT_GREY, Pt(12)),
    ("  Compare β-coefficients across models", False, LIGHT_GREY, Pt(12)),
    ("", False, WHITE, Pt(5)),
    ("Notebook 5 — Control Variables:", True, YELLOW, Pt(13)),
    ("  Partial out word length, frequency (log), POS,", False, LIGHT_GREY, Pt(12)),
    ("  sentence position — isolate pure surprisal effect.", False, LIGHT_GREY, Pt(12)),
    ("  Length and Frequency are huge confounds.", False, LIGHT_GREY, Pt(12)),
    ("", False, WHITE, Pt(5)),
    ("Notebook 6 — Syntactic Analysis (RQ3):", True, YELLOW, Pt(13)),
    ("  Locate specific garden-path sentences.", False, LIGHT_GREY, Pt(12)),
    ("  Plot word-by-word line chart: Surprisal vs RT.", False, LIGHT_GREY, Pt(12)),
    ("", False, WHITE, Pt(5)),
    ("Non-linear models: Random Forest / XGBoost", True, YELLOW, Pt(13)),
]
add_multiline(s, next_pts, Inches(0.55), Inches(1.95), Inches(7.1), Inches(4.6))

card(s, Inches(8.2), Inches(1.4), Inches(4.8), Inches(5.5))
add_text(s, "📚  References",
         Inches(8.35), Inches(1.5), Inches(4.5), Inches(0.4),
         font_size=Pt(15), bold=True, color=ACCENT_BLUE)
refs = [
    "Smith & Levy (2013).",
    "  Predictability effects on",
    "  reading time.",
    "",
    "Hale (2001).",
    "  A probabilistic Earley parser",
    "  as a psycholinguistic model.",
    "",
    "Levy (2008).",
    "  Expectation-based syntactic",
    "  comprehension.",
    "",
    "Demberg & Keller (2008).",
    "  Eye-tracking evidence for",
    "  sentence processing.",
    "",
    "Salazar et al. (2020).",
    "  Masked LM scoring.",
    "",
    "Kennedy et al. (2003).",
    "  Dundee eye-tracking corpus.",
]
add_multiline(s, refs, Inches(8.35), Inches(2.0), Inches(4.5), Inches(4.5),
              font_size=Pt(11.5), color=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
